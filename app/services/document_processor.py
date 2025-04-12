from typing import Optional
import docx
from pptx import Presentation
from pdfminer.high_level import extract_text
import os

def extract_text_from_file(file_path: str, file_type: str) -> Optional[str]:
    """
    Extract text content from different types of files
    """
    try:
        if file_type == 'docx':
            doc = docx.Document(file_path)
            return '\n'.join([paragraph.text for paragraph in doc.paragraphs])
        
        elif file_type == 'pptx':
            prs = Presentation(file_path)
            text = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text.append(shape.text)
            return '\n'.join(text)
        
        elif file_type == 'pdf':
            return extract_text(file_path)
        
        elif file_type == 'txt':
            with open(file_path, 'r', encoding='utf-8') as file:
                return file.read()
        
        else:
            raise ValueError(f"Unsupported file type: {file_type}")
    
    except Exception as e:
        print(f"Error extracting text from file: {str(e)}")
        return None
    finally:
        # Clean up the temporary file
        if os.path.exists(file_path):
            os.remove(file_path) 